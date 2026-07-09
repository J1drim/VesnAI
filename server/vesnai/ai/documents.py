"""Generate PDF, DOCX, and PPTX documents from structured content."""

from __future__ import annotations

import json
import re
import uuid
from dataclasses import dataclass, field
from io import BytesIO
from typing import Any

from vesnai.providers.base import AIProvider


@dataclass
class DocumentSection:
    heading: str
    body: str = ""
    bullets: list[str] = field(default_factory=list)


@dataclass
class DocumentSlide:
    title: str
    bullets: list[str] = field(default_factory=list)


@dataclass
class StructuredDocument:
    title: str
    sections: list[DocumentSection] = field(default_factory=list)
    slides: list[DocumentSlide] = field(default_factory=list)


class DocumentService:
    def __init__(self, reasoning: AIProvider) -> None:
        self.reasoning = reasoning

    def expand_outline(self, title: str, outline: str, fmt: str) -> StructuredDocument:
        prompt = (
            f"Expand this document outline into JSON for a {fmt} file.\n"
            f"Title: {title}\nOutline:\n{outline}\n\n"
            'Return ONLY JSON: {"title":"...","sections":[{"heading":"...","body":"...","bullets":[]}]} '
            'or for pptx {"title":"...","slides":[{"title":"...","bullets":[]}]}'
        )
        raw = self.reasoning.complete(prompt)
        data = _parse_json_object(raw)
        if fmt == "pptx":
            slides = [
                DocumentSlide(
                    title=str(s.get("title", "")),
                    bullets=[str(b) for b in s.get("bullets", [])],
                )
                for s in data.get("slides", [])
            ]
            return StructuredDocument(title=str(data.get("title", title)), slides=slides)
        sections = [
            DocumentSection(
                heading=str(s.get("heading", "")),
                body=str(s.get("body", "")),
                bullets=[str(b) for b in s.get("bullets", [])],
            )
            for s in data.get("sections", [])
        ]
        if not sections:
            sections = [DocumentSection(heading=title, body=outline)]
        return StructuredDocument(title=str(data.get("title", title)), sections=sections)

    def render(self, fmt: str, doc: StructuredDocument) -> tuple[bytes, str, str]:
        fmt = fmt.lower()
        if fmt == "pdf":
            return _render_pdf(doc), "application/pdf", ".pdf"
        if fmt == "docx":
            return (
                _render_docx(doc),
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                ".docx",
            )
        if fmt == "pptx":
            return (
                _render_pptx(doc),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                ".pptx",
            )
        raise ValueError(f"unsupported format: {fmt}")

    def generate(
        self, fmt: str, title: str, outline: str, *, session_stem: str = "chat"
    ) -> tuple[bytes, str, str]:
        structured = self.expand_outline(title, outline, fmt)
        data, mime, ext = self.render(fmt, structured)
        return data, mime, ext

    @staticmethod
    def attachment_path(session_stem: str, ext: str) -> str:
        stem = re.sub(r"[^a-zA-Z0-9]", "", session_stem)[:8] or "chat"
        return f"attachments/{stem}-doc-{uuid.uuid4().hex[:8]}{ext}"


def _parse_json_object(text: str) -> dict[str, Any]:
    text = text.strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?\s*", "", text)
        text = re.sub(r"\s*```$", "", text)
    start = text.find("{")
    end = text.rfind("}")
    if start >= 0 and end > start:
        text = text[start : end + 1]
    return json.loads(text)


def _render_pdf(doc: StructuredDocument) -> bytes:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import ListFlowable, ListItem, Paragraph, SimpleDocTemplate, Spacer

    buf = BytesIO()
    story = []
    styles = getSampleStyleSheet()
    pdf = SimpleDocTemplate(buf, pagesize=letter)
    story.append(Paragraph(doc.title, styles["Title"]))
    story.append(Spacer(1, 12))
    for section in doc.sections:
        if section.heading:
            story.append(Paragraph(section.heading, styles["Heading2"]))
        if section.body:
            story.append(Paragraph(section.body, styles["Normal"]))
        if section.bullets:
            items = [ListItem(Paragraph(b, styles["Normal"])) for b in section.bullets]
            story.append(ListFlowable(items, bulletType="bullet"))
        story.append(Spacer(1, 8))
    pdf.build(story)
    return buf.getvalue()


def _render_docx(doc: StructuredDocument) -> bytes:
    from docx import Document

    document = Document()
    document.add_heading(doc.title, level=0)
    for section in doc.sections:
        if section.heading:
            document.add_heading(section.heading, level=1)
        if section.body:
            document.add_paragraph(section.body)
        for bullet in section.bullets:
            document.add_paragraph(bullet, style="List Bullet")
    buf = BytesIO()
    document.save(buf)
    return buf.getvalue()


def _render_pptx(doc: StructuredDocument) -> bytes:
    from pptx import Presentation

    prs = Presentation()
    if not doc.slides:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = doc.title
    else:
        for i, s in enumerate(doc.slides):
            layout = prs.slide_layouts[1 if i else 0]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = s.title or doc.title
            body = slide.placeholders[1].text_frame
            body.clear()
            for bullet in s.bullets:
                p = body.add_paragraph()
                p.text = bullet
                p.level = 0
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
